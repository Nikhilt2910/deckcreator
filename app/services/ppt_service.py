from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
from uuid import uuid4

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from app.schemas.presentation import PresentationAnalysis


DEFAULT_TITLE_COLOR = RGBColor(28, 37, 44)
DEFAULT_BODY_COLOR = RGBColor(70, 78, 86)
DEFAULT_ACCENT_COLOR = RGBColor(13, 92, 99)


def generate_presentation(reference_path: Path, analysis: PresentationAnalysis, output_dir: Path) -> Path:
    presentation = _open_reference(reference_path)
    styles = _extract_template_style(presentation)
    reference_images = _extract_reference_images(presentation)

    if reference_path.suffix.lower() in {".pptx", ".potx"} and _can_populate_reference_deck(presentation):
        _populate_reference_deck(presentation, analysis)
    else:
        _clear_existing_slides(presentation)
        _build_generated_deck(presentation, analysis, styles, reference_images)

    output_path = output_dir / _build_output_name(analysis.title)
    presentation.save(output_path)
    return output_path


def _build_generated_deck(
    presentation: Presentation,
    analysis: PresentationAnalysis,
    styles: dict[str, object],
    reference_images: list[bytes],
) -> None:
    _add_title_slide(presentation, analysis, styles, reference_images)
    _add_kpi_slide(presentation, analysis, styles)
    _add_table_slide(presentation, "Channel Performance", analysis.channel_rows, styles, "Revenue", "ROI")
    _add_table_slide(presentation, "Regional Performance", analysis.region_rows, styles, "Revenue", "ROI")
    _add_table_slide(presentation, "Top Campaigns", analysis.top_campaign_rows, styles, "Revenue", "ROI")
    _add_summary_slide(presentation, analysis, styles, reference_images)
    _add_bullet_slide(presentation, "Key Insights", analysis.key_insights, styles)
    _add_bullet_slide(presentation, "Trends", analysis.trends, styles)
    _add_bullet_slide(presentation, "Risks", analysis.risks, styles)
    _add_sample_data_slide(presentation, analysis, styles)


def _open_reference(reference_path: Path) -> Presentation:
    if reference_path.suffix.lower() in {".pptx", ".potx"}:
        return Presentation(reference_path)
    return Presentation()


def _can_populate_reference_deck(presentation: Presentation) -> bool:
    text_shape_count = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text_shape_count += 1
    return len(presentation.slides) >= 6 and text_shape_count >= 12


def _populate_reference_deck(presentation: Presentation, analysis: PresentationAnalysis) -> None:
    slide_payloads = [
        [
            analysis.title,
            analysis.executive_summary[:140],
        ],
        [
            "Business Context",
            "Company",
            f"Campaign portfolio covers {analysis.kpis.get('campaign_count', '-')} marketing campaigns.",
            "Context",
            f"Revenue {analysis.kpis.get('total_revenue', '-')} from investment {analysis.kpis.get('total_investment', '-')}.",
            "Problem statement",
            analysis.executive_summary,
        ],
        [
            "Solution",
            "Performance focus",
            _first_or_default(analysis.key_insights, "The workbook highlights the best-performing segments and campaigns."),
        ],
        [
            "Challenges deep-dive",
            "Challenge 1",
            _first_or_default(analysis.risks, "Performance concentration can increase exposure to fewer segments."),
            "Challenge 2",
            _nth_or_default(analysis.risks, 1, _first_or_default(analysis.trends, "Efficiency varies significantly across channels.")),
            "Challenge 3",
            _nth_or_default(analysis.trends, 1, "Execution quality should remain consistent across regions."),
        ],
        [
            "Implementation",
        ],
        _build_timeline_payload(analysis),
        _build_team_payload(analysis),
        _build_impact_payload(analysis),
    ]

    for slide_index, slide in enumerate(presentation.slides):
        replacements = slide_payloads[slide_index] if slide_index < len(slide_payloads) else []
        _replace_placeholder_shapes(slide, replacements)


def _build_timeline_payload(analysis: PresentationAnalysis) -> list[str]:
    payload = ["Implementation"]
    campaign_rows = analysis.top_campaign_rows[:5]
    for index, row in enumerate(campaign_rows, start=1):
        payload.append(f"Top {index}")
        payload.append(f"{row.label}: {row.value_1} revenue at {row.value_2} ROI")
    return payload


def _build_team_payload(analysis: PresentationAnalysis) -> list[str]:
    payload = ["Leading segments"]
    source_rows = (analysis.region_rows + analysis.channel_rows)[:8]
    for row in source_rows:
        payload.append(row.label)
        payload.append(f"{row.value_1} revenue | {row.value_2} ROI")
    return payload


def _build_impact_payload(analysis: PresentationAnalysis) -> list[str]:
    return [
        "Impact",
        analysis.kpis.get("average_roi", "-"),
        analysis.kpis.get("total_revenue", "-"),
        analysis.kpis.get("campaign_count", "-"),
    ]


def _replace_placeholder_shapes(slide, replacements: list[str]) -> None:
    placeholder_shapes = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER
    ]
    for index, shape in enumerate(placeholder_shapes):
        replacement = replacements[index] if index < len(replacements) else ""
        _set_shape_text(shape, replacement)


def _set_shape_text(shape, text: str) -> None:
    if not shape.has_text_frame:
        return

    style = _capture_shape_style(shape)
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    paragraph = text_frame.paragraphs[0]
    paragraph.text = _fit_text_to_shape(text, shape, style["font_size"])
    paragraph.alignment = style["alignment"]
    paragraph.font.name = style["font_name"]
    paragraph.font.size = style["font_size"]
    paragraph.font.bold = style["bold"]
    if style["color"] is not None:
        paragraph.font.color.rgb = style["color"]


def _clear_existing_slides(presentation: Presentation) -> None:
    slide_ids = list(presentation.slides._sldIdLst)
    for slide_id in slide_ids:
        relationship_id = slide_id.rId
        presentation.part.drop_rel(relationship_id)
        presentation.slides._sldIdLst.remove(slide_id)


def _extract_reference_images(presentation: Presentation) -> list[bytes]:
    images: list[bytes] = []
    seen_hashes: set[int] = set()
    for slide in presentation.slides:
        for shape in slide.shapes:
            images.extend(_extract_images_from_shape(shape, seen_hashes))
    return images[:6]


def _extract_images_from_shape(shape, seen_hashes: set[int]) -> list[bytes]:
    images: list[bytes] = []
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        blob = shape.image.blob
        blob_hash = hash(blob)
        if blob_hash not in seen_hashes:
            seen_hashes.add(blob_hash)
            images.append(blob)
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for group_shape in shape.shapes:
            images.extend(_extract_images_from_shape(group_shape, seen_hashes))
    return images


def _extract_template_style(presentation: Presentation) -> dict[str, object]:
    title_font_name = "Aptos Display"
    body_font_name = "Aptos"
    title_size = Pt(28)
    body_size = Pt(18)

    if presentation.slides:
        first_slide = presentation.slides[0]
        if first_slide.shapes.title and first_slide.shapes.title.has_text_frame:
            title_paragraph = first_slide.shapes.title.text_frame.paragraphs[0]
            if title_paragraph.font.name:
                title_font_name = title_paragraph.font.name
            if title_paragraph.font.size:
                title_size = title_paragraph.font.size

        for shape in first_slide.shapes:
            if shape.has_text_frame and shape != first_slide.shapes.title:
                paragraph = shape.text_frame.paragraphs[0]
                if paragraph.font.name:
                    body_font_name = paragraph.font.name
                if paragraph.font.size:
                    body_size = paragraph.font.size
                break

    return {
        "title_font_name": title_font_name,
        "body_font_name": body_font_name,
        "title_size": title_size,
        "body_size": body_size,
    }


def _add_title_slide(
    presentation: Presentation,
    analysis: PresentationAnalysis,
    styles: dict[str, object],
    reference_images: list[bytes],
) -> None:
    slide = presentation.slides.add_slide(_get_layout(presentation, "blank"))

    title_width = Inches(5.3 if reference_images else 8.0)
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), title_width, Inches(1.6))
    title.text_frame.text = analysis.title
    _format_text_frame(title.text_frame, str(styles["title_font_name"]), Pt(max(int(styles["title_size"].pt), 28)), DEFAULT_TITLE_COLOR, True)

    subtitle = slide.shapes.add_textbox(Inches(0.75), Inches(2.05), title_width, Inches(1.0))
    subtitle.text_frame.text = "Automated marketing performance review generated from uploaded workbook data"
    _format_text_frame(subtitle.text_frame, str(styles["body_font_name"]), Pt(18), DEFAULT_ACCENT_COLOR, False)

    if reference_images:
        slide.shapes.add_picture(BytesIO(reference_images[0]), Inches(6.1), Inches(0.7), width=Inches(3.0), height=Inches(4.6))


def _add_kpi_slide(presentation: Presentation, analysis: PresentationAnalysis, styles: dict[str, object]) -> None:
    slide = presentation.slides.add_slide(_get_layout(presentation, "blank"))
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(4.8), Inches(0.8))
    title.text_frame.text = "Performance Overview"
    _format_text_frame(title.text_frame, str(styles["title_font_name"]), Pt(24), DEFAULT_TITLE_COLOR, True)

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(8.1), Inches(4.8))
    metrics = [
        f"Campaigns: {analysis.kpis.get('campaign_count', '-')}",
        f"Total revenue: {analysis.kpis.get('total_revenue', '-')}",
        f"Total investment: {analysis.kpis.get('total_investment', '-')}",
        f"Average ROI: {analysis.kpis.get('average_roi', '-')}",
    ]
    for index, metric in enumerate(metrics):
        paragraph = body.text_frame.paragraphs[0] if index == 0 else body.text_frame.add_paragraph()
        paragraph.text = metric
        _format_paragraph(paragraph, str(styles["body_font_name"]), Pt(22), DEFAULT_BODY_COLOR, index == 0)


def _add_table_slide(
    presentation: Presentation,
    slide_title: str,
    rows,
    styles: dict[str, object],
    col_two_title: str,
    col_three_title: str,
) -> None:
    if not rows:
        return

    slide = presentation.slides.add_slide(_get_layout(presentation, "blank"))
    title = slide.shapes.add_textbox(Inches(0.65), Inches(0.45), Inches(5.0), Inches(0.8))
    title.text_frame.text = slide_title
    _format_text_frame(title.text_frame, str(styles["title_font_name"]), Pt(24), DEFAULT_TITLE_COLOR, True)

    visible_rows = rows[:6]
    table_shape = slide.shapes.add_table(len(visible_rows) + 1, 3, Inches(0.55), Inches(1.35), Inches(9.0), Inches(4.85))
    table = table_shape.table
    table.columns[0].width = Inches(4.35)
    table.columns[1].width = Inches(2.15)
    table.columns[2].width = Inches(2.5)

    headers = ["Segment", col_two_title, col_three_title]
    for index, header in enumerate(headers):
        table.cell(0, index).text = header
        _format_cell(table.cell(0, index), str(styles["body_font_name"]), Pt(15), DEFAULT_TITLE_COLOR, True)

    for row_index, row in enumerate(visible_rows, start=1):
        values = [row.label, row.value_1, row.value_2]
        for col_index, value in enumerate(values):
            table.cell(row_index, col_index).text = value
            _format_cell(table.cell(row_index, col_index), str(styles["body_font_name"]), Pt(14), DEFAULT_BODY_COLOR, False)


def _add_summary_slide(
    presentation: Presentation,
    analysis: PresentationAnalysis,
    styles: dict[str, object],
    reference_images: list[bytes],
) -> None:
    slide = presentation.slides.add_slide(_get_layout(presentation, "blank"))
    title = slide.shapes.add_textbox(Inches(0.65), Inches(0.5), Inches(4.6), Inches(0.8))
    title.text_frame.text = "Executive Summary"
    _format_text_frame(title.text_frame, str(styles["title_font_name"]), Pt(24), DEFAULT_TITLE_COLOR, True)

    body_width = Inches(5.9 if len(reference_images) > 1 else 8.3)
    body = slide.shapes.add_textbox(Inches(0.75), Inches(1.35), body_width, Inches(4.85))
    body.text_frame.word_wrap = True
    paragraph = body.text_frame.paragraphs[0]
    paragraph.text = analysis.executive_summary
    _format_paragraph(paragraph, str(styles["body_font_name"]), Pt(17), DEFAULT_BODY_COLOR, False)

    if len(reference_images) > 1:
        _add_reference_image_strip(slide, reference_images[1:4], Inches(6.85), Inches(1.55), Inches(1.95), Inches(1.2), gap=0.16)


def _add_bullet_slide(presentation: Presentation, slide_title: str, items: list[str], styles: dict[str, object]) -> None:
    if not items:
        return

    slide = presentation.slides.add_slide(_get_layout(presentation, "blank"))
    title = slide.shapes.add_textbox(Inches(0.65), Inches(0.5), Inches(4.2), Inches(0.8))
    title.text_frame.text = slide_title
    _format_text_frame(title.text_frame, str(styles["title_font_name"]), Pt(24), DEFAULT_TITLE_COLOR, True)

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(8.2), Inches(4.9))
    body.text_frame.word_wrap = True
    for index, item in enumerate(items[:5]):
        paragraph = body.text_frame.paragraphs[0] if index == 0 else body.text_frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        _format_paragraph(paragraph, str(styles["body_font_name"]), Pt(16), DEFAULT_BODY_COLOR, False)


def _add_sample_data_slide(presentation: Presentation, analysis: PresentationAnalysis, styles: dict[str, object]) -> None:
    if not analysis.sample_rows:
        return

    slide = presentation.slides.add_slide(_get_layout(presentation, "blank"))
    title = slide.shapes.add_textbox(Inches(0.65), Inches(0.45), Inches(4.6), Inches(0.8))
    title.text_frame.text = "Workbook Sample"
    _format_text_frame(title.text_frame, str(styles["title_font_name"]), Pt(24), DEFAULT_TITLE_COLOR, True)

    columns = list(analysis.sample_rows[0].keys())[:4]
    row_count = min(len(analysis.sample_rows), 5)
    table_shape = slide.shapes.add_table(row_count + 1, len(columns), Inches(0.45), Inches(1.35), Inches(9.1), Inches(4.95))
    table = table_shape.table

    for col_index, header in enumerate(columns):
        table.cell(0, col_index).text = header
        _format_cell(table.cell(0, col_index), str(styles["body_font_name"]), Pt(13), DEFAULT_TITLE_COLOR, True)

    for row_index in range(row_count):
        for col_index, header in enumerate(columns):
            table.cell(row_index + 1, col_index).text = analysis.sample_rows[row_index].get(header, "")
            _format_cell(table.cell(row_index + 1, col_index), str(styles["body_font_name"]), Pt(12), DEFAULT_BODY_COLOR, False)


def _get_layout(presentation: Presentation, kind: str):
    match kind:
        case "blank":
            return _find_layout_by_name(presentation, ["blank"]) or presentation.slide_layouts[0]
        case _:
            return presentation.slide_layouts[0]


def _find_layout_by_name(presentation: Presentation, candidates: list[str]):
    normalized = {candidate.lower() for candidate in candidates}
    for layout in presentation.slide_layouts:
        if layout.name.lower() in normalized:
            return layout
    return None


def _format_text_frame(text_frame, font_name: str, font_size, color: RGBColor, bold: bool) -> None:
    for paragraph in text_frame.paragraphs:
        _format_paragraph(paragraph, font_name, font_size, color, bold)


def _format_paragraph(paragraph, font_name: str, font_size, color: RGBColor, bold: bool) -> None:
    paragraph.font.name = font_name
    paragraph.font.size = font_size
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.alignment = PP_ALIGN.LEFT


def _format_cell(cell, font_name: str, font_size, color: RGBColor, bold: bool) -> None:
    cell.text_frame.word_wrap = True
    cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for paragraph in cell.text_frame.paragraphs:
        _format_paragraph(paragraph, font_name, font_size, color, bold)


def _add_reference_image_strip(slide, images: list[bytes], start_left, top, width, height, gap: float) -> None:
    current_left = start_left
    for blob in images:
        slide.shapes.add_picture(BytesIO(blob), current_left, top, width=width, height=height)
        current_left += width + Inches(gap)


def _first_or_default(items: list[str], default: str) -> str:
    return items[0] if items else default


def _nth_or_default(items: list[str], index: int, default: str) -> str:
    return items[index] if len(items) > index else default


def _capture_shape_style(shape) -> dict[str, object]:
    default_style = {
        "font_name": "Aptos",
        "font_size": Pt(18),
        "bold": False,
        "alignment": PP_ALIGN.LEFT,
        "color": None,
    }

    if not shape.has_text_frame or not shape.text_frame.paragraphs:
        return default_style

    paragraph = shape.text_frame.paragraphs[0]
    font = paragraph.font
    color = None
    if font.color is not None and getattr(font.color, "rgb", None) is not None:
        color = font.color.rgb

    return {
        "font_name": font.name or default_style["font_name"],
        "font_size": font.size or default_style["font_size"],
        "bold": bool(font.bold) if font.bold is not None else default_style["bold"],
        "alignment": paragraph.alignment or default_style["alignment"],
        "color": color,
    }


def _fit_text_to_shape(text: str, shape, font_size) -> str:
    cleaned = " ".join(text.split())
    if not cleaned:
        return ""

    font_points = max(font_size.pt if font_size is not None else 18, 10)
    width_points = max(shape.width / 12700, 40)
    height_points = max(shape.height / 12700, 20)

    chars_per_line = max(int(width_points / (font_points * 0.56)), 8)
    line_count = max(int(height_points / (font_points * 1.5)), 1)
    max_chars = max(chars_per_line * line_count, 12)

    if len(cleaned) <= max_chars:
        return cleaned

    truncated = cleaned[: max_chars - 1].rstrip()
    if " " in truncated:
        truncated = truncated.rsplit(" ", 1)[0]
    return f"{truncated}..."


def _build_output_name(title: str) -> str:
    slug = "".join(char.lower() if char.isalnum() else "-" for char in title).strip("-")
    clean_slug = "-".join(part for part in slug.split("-") if part)[:40] or "presentation"
    timestamp = datetime.now(timezone.utc).strftime("%Y%m%d%H%M%S")
    return f"{clean_slug}-{timestamp}-{uuid4().hex[:8]}.pptx"
