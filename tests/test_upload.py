from io import BytesIO
from pathlib import Path
import sys
import unittest
import json
import tempfile
from unittest.mock import patch

import pandas as pd
from fastapi.testclient import TestClient
from pptx import Presentation
from pypdf import PdfWriter

PROJECT_ROOT = Path(__file__).resolve().parent.parent
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from app.main import app
from app.services.engineering_agent_service import _try_generate_literal_text_resolution
from app.services.patch_service import apply_unified_diff, validate_unified_diff
from app.services.review_token_service import build_review_token
from app.schemas.ticket import TicketAutomationResult, TicketResolution, TicketReviewOutcome


class UploadApiTestCase(unittest.TestCase):
    def setUp(self) -> None:
        self.client = TestClient(app)

    def test_root_endpoint(self) -> None:
        response = self.client.get("/")

        self.assertEqual(response.status_code, 200)
        self.assertIn("text/html", response.headers["content-type"])
        self.assertIn("Generate Presentation", response.text)
        self.assertIn("Case Study", response.text)

    def test_upload_endpoint_saves_excel_and_template(self) -> None:
        excel_bytes = self._build_excel_file()
        pdf_bytes = self._build_pdf_reference()

        response = self.client.post(
            "/upload",
            files={
                "excel_file": (
                    "sample.xlsx",
                    excel_bytes.getvalue(),
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                ),
                "reference_file": (
                    "reference.pdf",
                    pdf_bytes.getvalue(),
                    "application/pdf",
                ),
            },
        )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertEqual(payload["message"], "Files uploaded successfully.")
        self.assertTrue(Path(payload["excel_file"]["saved_path"]).exists())
        self.assertTrue(Path(payload["reference_file"]["saved_path"]).exists())

    def test_generate_presentation_returns_downloadable_pptx(self) -> None:
        ppt_bytes = self._build_ppt_template()
        analysis_payload = {
            "title": "Q2 Business Review",
            "executive_summary": "Revenue grew ahead of plan while margin pressure increased in two regions.",
            "key_insights": [
                "North America led revenue growth.",
                "Enterprise accounts expanded faster than SMB.",
            ],
            "trends": [
                "Average deal size increased quarter over quarter.",
                "Renewal rates improved in strategic segments.",
            ],
            "risks": [
                "Margin compression is emerging in EMEA.",
            ],
        }

        response = self.client.post(
            "/presentations/generate",
            files={
                "ppt_template": (
                    "template.pptx",
                    ppt_bytes.getvalue(),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                ),
            },
            data={"analysis_json": json.dumps(analysis_payload)},
        )

        self.assertEqual(response.status_code, 200)
        self.assertIn(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            response.headers["content-type"],
        )
        self.assertTrue(response.content.startswith(b"PK"))

    @patch("app.services.report_service.analyze_data")
    def test_report_generation_endpoint_returns_downloadable_pptx(self, mock_analyze_data) -> None:
        excel_bytes = self._build_excel_file()
        ppt_bytes = self._build_ppt_template()
        mock_analyze_data.return_value = {
            "key_insights": ["Revenue outperformed expectations."],
            "trends": ["Demand is shifting toward larger enterprise deals."],
            "risks": ["Margin pressure is increasing in one region."],
            "executive_summary": "The quarter showed healthy growth with targeted margin risks.",
        }

        response = self.client.post(
            "/reports/generate",
            files={
                "excel_file": (
                    "sample.xlsx",
                    excel_bytes.getvalue(),
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                ),
                "reference_file": (
                    "template.pptx",
                    ppt_bytes.getvalue(),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                ),
            },
        )

        self.assertEqual(response.status_code, 200)
        self.assertIn(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            response.headers["content-type"],
        )
        self.assertTrue(response.content.startswith(b"PK"))

    @patch("app.services.report_service.analyze_data")
    def test_report_generation_accepts_pdf_reference(self, mock_analyze_data) -> None:
        excel_bytes = self._build_excel_file()
        pdf_bytes = self._build_pdf_reference()
        mock_analyze_data.return_value = {
            "key_insights": ["Email delivered strong ROI."],
            "trends": ["Revenue concentration is highest in one segment."],
            "risks": ["One channel is underperforming on efficiency."],
            "executive_summary": "The dataset shows healthy returns with a few concentrated risks.",
        }

        response = self.client.post(
            "/reports/generate",
            files={
                "excel_file": (
                    "sample.xlsx",
                    excel_bytes.getvalue(),
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                ),
                "reference_file": (
                    "reference.pdf",
                    pdf_bytes.getvalue(),
                    "application/pdf",
                ),
            },
        )

        self.assertEqual(response.status_code, 200)
        self.assertTrue(response.content.startswith(b"PK"))

    @patch("app.services.report_service.analyze_data")
    def test_report_generation_accepts_built_in_template(self, mock_analyze_data) -> None:
        excel_bytes = self._build_excel_file()
        mock_analyze_data.return_value = {
            "key_insights": ["Email delivered strong ROI."],
            "trends": ["Revenue concentration is highest in one segment."],
            "risks": ["One channel is underperforming on efficiency."],
            "executive_summary": "The dataset shows healthy returns with a few concentrated risks.",
        }

        response = self.client.post(
            "/reports/generate",
            files={
                "excel_file": (
                    "sample.xlsx",
                    excel_bytes.getvalue(),
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                ),
            },
            data={"built_in_template": "case-study"},
        )

        self.assertEqual(response.status_code, 200)
        self.assertTrue(response.content.startswith(b"PK"))

    def test_ticket_endpoint_stores_ticket_locally(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            tickets_path = Path(temp_dir) / "tickets.json"
            with patch("app.services.ticket_service.TICKETS_FILE", tickets_path), patch(
                "app.services.ticket_service._create_jira_issue",
                return_value=None,
            ), patch(
                "app.services.ticket_service.generate_ticket_resolution",
                return_value=TicketResolution(
                    files=["app/services/ticket_service.py"],
                    patch="--- a/app/services/ticket_service.py\n+++ b/app/services/ticket_service.py",
                    explanation="Update ticket persistence.",
                    generated_at="2026-04-18T00:00:00Z",
                ),
            ), patch(
                "app.services.ticket_service.send_ticket_review_email",
                return_value=(True, None),
            ):
                response = self.client.post(
                    "/ticket",
                    json={
                        "type": "feature",
                        "description": "Add slide previews to the built-in template gallery.",
                    },
                )

            self.assertEqual(response.status_code, 200)
            payload = response.json()
            self.assertEqual(payload["type"], "feature")
            self.assertFalse(payload["jira_synced"])
            self.assertIsNone(payload["jira_issue_key"])
            self.assertEqual(payload["resolution"]["files"], ["app/services/ticket_service.py"])
            self.assertEqual(payload["status"], "pending")
            self.assertTrue(payload["email_sent"])
            self.assertTrue(tickets_path.exists())

            stored_payload = json.loads(tickets_path.read_text(encoding="utf-8"))
            self.assertEqual(len(stored_payload), 1)
            self.assertEqual(stored_payload[0]["description"], payload["description"])

    def test_ticket_endpoint_skips_email_when_no_valid_resolution_exists(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            tickets_path = Path(temp_dir) / "tickets.json"
            with patch("app.services.ticket_service.TICKETS_FILE", tickets_path), patch(
                "app.services.ticket_service._create_jira_issue",
                return_value=None,
            ), patch(
                "app.services.ticket_service.generate_ticket_resolution",
                side_effect=ValueError("non-applyable patch"),
            ), patch(
                "app.services.ticket_service.send_ticket_review_email",
            ) as mock_email:
                response = self.client.post(
                    "/ticket",
                    json={
                        "type": "bug",
                        "description": "Remove a small piece of text from the UI.",
                    },
                )

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertIsNone(payload["resolution"])
        self.assertFalse(payload["email_sent"])
        self.assertIn("No valid resolution", payload["email_error"])
        mock_email.assert_not_called()

    def test_ticket_resolution_regeneration_updates_existing_ticket(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            tickets_path = Path(temp_dir) / "tickets.json"
            seeded_ticket = [{
                "id": "ticket123",
                "type": "bug",
                "description": "Fix the upload status behavior.",
                "created_at": "2026-04-18T00:00:00Z",
                "jira_synced": False,
                "jira_issue_key": None,
                "resolution": None,
            }]
            tickets_path.write_text(json.dumps(seeded_ticket), encoding="utf-8")
            with patch("app.services.ticket_service.TICKETS_FILE", tickets_path), patch(
                "app.services.ticket_service.generate_ticket_resolution",
                return_value=TicketResolution(
                    files=["templates/index.html", "static/styles.css"],
                    patch="--- a/templates/index.html\n+++ b/templates/index.html",
                    explanation="Align feedback status behavior.",
                    generated_at="2026-04-18T00:00:00Z",
                ),
            ):
                response = self.client.post("/ticket/ticket123/resolution")

            self.assertEqual(response.status_code, 200)
            payload = response.json()
            self.assertEqual(payload["id"], "ticket123")
            self.assertEqual(payload["resolution"]["files"], ["templates/index.html", "static/styles.css"])

    def test_ticket_review_page_renders_resolution(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            tickets_path = Path(temp_dir) / "tickets.json"
            tickets_path.write_text(json.dumps([{
                "id": "ticket999",
                "type": "feature",
                "description": "Add review page.",
                "created_at": "2026-04-18T00:00:00Z",
                "jira_synced": False,
                "jira_issue_key": None,
                "resolution": {
                    "files": ["templates/ticket_review.html"],
                    "patch": "--- a/templates/ticket_review.html\n+++ b/templates/ticket_review.html",
                    "explanation": "Render a review page.",
                    "generated_at": "2026-04-18T00:00:00Z",
                },
                "status": "pending",
                "developer_email": "nikhil.t2910@gmail.com",
                "review_url": "http://127.0.0.1:8000/ticket/ticket999/review",
                "email_sent": False,
                "email_error": "SMTP is not configured.",
                "review_outcome": None,
            }]), encoding="utf-8")
            with patch("app.services.ticket_service.TICKETS_FILE", tickets_path):
                token = build_review_token("ticket999")
                response = self.client.get(f"/ticket/ticket999/review?token={token}")

        self.assertEqual(response.status_code, 200)
        self.assertIn("Developer Review", response.text)
        self.assertIn("Add review page.", response.text)

    def test_ticket_approve_applies_patch_result(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            tickets_path = Path(temp_dir) / "tickets.json"
            tickets_path.write_text(json.dumps([{
                "id": "ticket777",
                "type": "bug",
                "description": "Fix review button state.",
                "created_at": "2026-04-18T00:00:00Z",
                "jira_synced": False,
                "jira_issue_key": None,
                "resolution": {
                    "files": ["templates/ticket_review.html"],
                    "patch": "--- a/templates/ticket_review.html\n+++ b/templates/ticket_review.html",
                    "explanation": "Adjust the buttons.",
                    "generated_at": "2026-04-18T00:00:00Z",
                },
                "status": "pending",
                "developer_email": "nikhil.t2910@gmail.com",
                "review_url": "http://127.0.0.1:8000/ticket/ticket777/review",
                "email_sent": True,
                "email_error": None,
                "review_outcome": None,
            }]), encoding="utf-8")
            with patch("app.services.ticket_service.TICKETS_FILE", tickets_path), patch(
                "app.services.ticket_service.apply_unified_diff",
                return_value=TicketReviewOutcome(
                    applied=True,
                    message="Patch applied successfully.",
                    applied_at="2026-04-18T00:00:00Z",
                ),
            ), patch(
                "app.services.ticket_service.run_post_approval_pipeline",
                return_value=TicketAutomationResult(
                    patch_applied=True,
                    tests_passed=True,
                    pushed=True,
                    branch="ticket/ticket777",
                    commit_sha="abc123",
                    message="Patch applied, tests passed, and changes were pushed to GitHub.",
                    completed_at="2026-04-18T00:00:00Z",
                ),
            ):
                token = build_review_token("ticket777")
                response = self.client.post(f"/ticket/ticket777/approve?token={token}")

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertEqual(payload["status"], "approved")
        self.assertTrue(payload["review_outcome"]["applied"])
        self.assertTrue(payload["automation_result"]["pushed"])

    def test_invalid_review_token_is_rejected(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            tickets_path = Path(temp_dir) / "tickets.json"
            tickets_path.write_text(json.dumps([{
                "id": "ticket401",
                "type": "bug",
                "description": "Reject invalid token access.",
                "created_at": "2026-04-18T00:00:00Z",
                "jira_synced": False,
                "jira_issue_key": None,
                "resolution": None,
                "status": "pending",
                "developer_email": "nikhil.t2910@gmail.com",
                "review_url": "http://127.0.0.1:8000/ticket/ticket401/review",
                "email_sent": False,
                "email_error": None,
                "review_outcome": None,
            }]), encoding="utf-8")
            with patch("app.services.ticket_service.TICKETS_FILE", tickets_path):
                response = self.client.get("/ticket/ticket401/review?token=bad-token")

        self.assertEqual(response.status_code, 403)

    def test_patch_service_uses_configured_git_executable(self) -> None:
        with patch.dict("os.environ", {"GIT_EXECUTABLE": r"C:\Program Files\Git\cmd\git.exe"}), patch(
            "app.services.patch_service.Path.exists",
            return_value=True,
        ), patch(
            "app.services.patch_service.subprocess.run",
        ) as mock_run:
            mock_run.return_value.returncode = 0
            mock_run.return_value.stdout = ""
            mock_run.return_value.stderr = ""

            result = apply_unified_diff("--- a/README.md\n+++ b/README.md\n")

        self.assertTrue(result.applied)
        command = mock_run.call_args.args[0]
        self.assertEqual(command[0], r"C:\Program Files\Git\cmd\git.exe")

    def test_patch_service_rejects_placeholder_patch(self) -> None:
        result = apply_unified_diff(
            "--- a/templates/index.html\n+++ b/templates/index.html\n@@ ... @@\n-foo\n+bar\n"
        )

        self.assertFalse(result.applied)
        self.assertIn("placeholder diff", result.message)

    def test_validate_unified_diff_rejects_placeholder_patch(self) -> None:
        is_valid, message = validate_unified_diff(
            "--- a/templates/index.html\n+++ b/templates/index.html\n@@ ... @@\n-foo\n+bar\n"
        )

        self.assertFalse(is_valid)
        self.assertIn("placeholder diff", message)

    def test_literal_text_resolution_handles_simple_remove_request(self) -> None:
        resolution = _try_generate_literal_text_resolution(
            'please remove "Start from a proven deck structure" in the page.'
        )

        self.assertIsNotNone(resolution)
        self.assertEqual(resolution.files, ["templates/index.html"])
        self.assertIn('-                <h2>Start from a proven deck structure</h2>', resolution.patch)
        self.assertIn('+                <h2></h2>', resolution.patch)

    def test_literal_text_resolution_supports_single_quoted_text(self) -> None:
        resolution = _try_generate_literal_text_resolution(
            "please remove 'potx' under the Inputs block"
        )

        self.assertIsNotNone(resolution)
        self.assertEqual(resolution.files, ["templates/index.html"])
        self.assertIn("-                <li>Reference file: `.pptx`, `.potx`, `.pdf`</li>", resolution.patch)
        self.assertIn("+                <li>Reference file: `.pptx`, `.pdf`</li>", resolution.patch)
        self.assertIn('-                <input id="reference_file" name="reference_file" type="file" accept=".pptx,.potx,.pdf" required>', resolution.patch)
        self.assertIn('+                <input id="reference_file" name="reference_file" type="file" accept=".pptx,.pdf" required>', resolution.patch)

    @staticmethod
    def _build_excel_file() -> BytesIO:
        buffer = BytesIO()
        frame = pd.DataFrame({"region": ["West", "East"], "sales": [120, 150]})
        with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
            frame.to_excel(writer, index=False, sheet_name="Summary")
        buffer.seek(0)
        return buffer

    @staticmethod
    def _build_ppt_template() -> BytesIO:
        buffer = BytesIO()
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[0])
        presentation.save(buffer)
        buffer.seek(0)
        return buffer

    @staticmethod
    def _build_pdf_reference() -> BytesIO:
        buffer = BytesIO()
        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        writer.write(buffer)
        buffer.seek(0)
        return buffer


if __name__ == "__main__":
    unittest.main()
